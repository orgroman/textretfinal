from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image, ImageDraw
import tempfile
import os
import math

# --- PREMIUM THEME (Modern Dark) ---
THEME = {
    'bg_dark': (20, 20, 20),           # Almost Black
    'bg_card': (35, 35, 35),           # Dark Gray for Cards
    'text_primary': (255, 255, 255),   # White
    'text_secondary': (180, 180, 180), # Light Gray
    'accent_gradient_start': (0, 120, 215),  # Microsoft Blue
    'accent_gradient_end': (0, 255, 255),    # Cyan
    'accent_warning': (255, 100, 100)        # Soft Red for warnings/failures
}

def create_gradient_img(width, height, start_color, end_color):
    img = Image.new('RGB', (width, height))
    for x in range(width):
        for y in range(height):
            ratio = (x/width + y/height) / 2
            r = int(start_color[0] * (1 - ratio) + end_color[0] * ratio)
            g = int(start_color[1] * (1 - ratio) + end_color[1] * ratio)
            b = int(start_color[2] * (1 - ratio) + end_color[2] * ratio)
            img.putpixel((x, y), (r, g, b))
    return img

def add_gradient_shape(slide, left, top, width, height, color1=None, color2=None):
    if not color1: color1 = THEME['accent_gradient_start']
    if not color2: color2 = THEME['accent_gradient_end']
    
    w_px, h_px = int(width.inches * 96), int(height.inches * 96)
    if w_px < 1 or h_px < 1: return
    
    img = create_gradient_img(w_px, h_px, color1, color2)
    with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as tmp:
        img.save(tmp.name)
        tmp_path = tmp.name
    pic = slide.shapes.add_picture(tmp_path, left, top, width, height)
    if os.path.exists(tmp_path): os.unlink(tmp_path)
    return pic

def setup_slide(prs, layout_index=6):
    slide = prs.slides.add_slide(prs.slide_layouts[layout_index])
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(*THEME['bg_dark'])
    background.line.fill.background()
    add_gradient_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.15))
    return slide

def add_title(slide, text, subtitle=None):
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
    p = title_box.text_frame.paragraphs[0]
    p.text = text.upper()
    p.font.name = 'Segoe UI Black'
    p.font.size = Pt(36)
    p.font.color.rgb = RGBColor(*THEME['text_primary'])
    add_gradient_shape(slide, Inches(0.5), Inches(1.3), Inches(2), Inches(0.05))
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(12), Inches(0.8))
        p = sub_box.text_frame.paragraphs[0]
        p.text = subtitle
        p.font.name = 'Segoe UI Light'
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(*THEME['text_secondary'])

def add_card_content(slide, title, bullets, left=Inches(0.5), top=Inches(2.2), width=Inches(12.3)):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, Inches(4.8))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(*THEME['bg_card'])
    card.line.fill.background()
    
    if title:
        tb = slide.shapes.add_textbox(left + Inches(0.3), top + Inches(0.2), width - Inches(0.6), Inches(0.8))
        p = tb.text_frame.paragraphs[0]
        p.text = title
        p.font.name = 'Segoe UI Semibold'
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(*THEME['accent_gradient_end'])
        
    border_offset = Inches(0.8) if title else Inches(0.3)
    tb = slide.shapes.add_textbox(left + Inches(0.3), top + border_offset, width - Inches(0.6), Inches(3.8))
    tf = tb.text_frame
    tf.word_wrap = True
    for line in bullets:
        p = tf.add_paragraph()
        p.text = line
        p.font.name = 'Segoe UI'
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(*THEME['text_primary'])
        p.space_after = Pt(12)

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # 1. TITLE
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = RGBColor(*THEME['bg_dark'])
    
    tb = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.3), Inches(2.5))
    p = tb.text_frame.paragraphs[0]
    p.text = "ROBUST04 RETRIEVAL:\nZERO-SHOT PIPELINE"
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Segoe UI Black'
    p.font.size = Pt(54)
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    tb = slide.shapes.add_textbox(Inches(2), Inches(5), Inches(9.3), Inches(1))
    p = tb.text_frame.paragraphs[0]
    p.text = "Maximizing Score with Limited Labels"
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Segoe UI Light'
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(0, 255, 255)

    # 2. GOAL
    slide = setup_slide(prs)
    add_title(slide, "The Goal: Zero-Shot Robustness", "Constraint: 50 Judged Queries")
    add_card_content(slide, "Why not train?", [
        "• Challenge: Only 50 labeled topics available.",
        "• Risk: Training complex models guarantees overfitting.",
        "• Strategy: ZERO-SHOT Assembly.",
        "   - Use pre-trained components (Pyserini, BGE, MonoT5).",
        "   - Optimize generic hyperparameters, not model weights."
    ])

    # 3. THREE RUN METHODS
    slide = setup_slide(prs)
    add_title(slide, "Methodology: Three Runs", "Isolating Improvements")
    add_card_content(slide, None, [
        "Run 1: Fusion Baseline",
        "   - Combines BM25(RM3) + SPLADE + Dense(HyDE).",
        "Run 2: Structural Validation (SDM)",
        "   - Replaces BM25 with Sequential Dependence Model (SDM).",
        "Run 3: Passage-Based Reranking",
        "   - Rescores Run 1 documents using MonoT5-3B (MaxP windowing)."
    ])

    # 4. METHOD 1: SPLADE (Common to all)
    slide = setup_slide(prs)
    add_title(slide, "Component: Learned Sparse (SPLADE)", "Solving Vocabulary Mismatch")
    add_card_content(slide, "The 'Semantic Anchor'", [
        "• Expands query with latent terms (e.g., 'car' -> 'vehicle').",
        "• Efficiency: Uses inverted index (fast) but finds synonyms.",
        "• Role in runs: Provides high-recall candidates for all 3 runs."
    ])

    # 5. METHOD 1: DENSE + HYDE
    slide = setup_slide(prs)
    add_title(slide, "Component: Hybrid Dense (HyDE)", "Zero-Shot Semantic Matching")
    add_card_content(slide, "The Hallucination Strategy", [
        "• Problem: Generic dense models drift on specific domains.",
        "• Solution (HyDE):",
        "   1. LLM generates a 'fake' relevant document.",
        "   2. BGE model embeds the fake doc.",
        "   3. Retrieves real docs matching the semantic intent."
    ])

    # 6. HYDE EXAMPLE
    slide = setup_slide(prs)
    add_title(slide, "HyDE in Action", "Success vs. Drift")
    add_card_content(slide, "Success: 'Hubble Achievements'", [
        "• Query: 'Hubble Telescope Achievements'",
        "• Hallucinated: '...supernova... exoplanets... dark energy...'",
        "• Result: Dense model finds matches missing keywords."
    ], width=Inches(5.8), top=Inches(2.5))
    
    add_card_content(slide, "Failure: 'Radio Waves Cancer'", [
        "• Query: 'Radio Waves Brain Cancer'",
        "• Hallucinated: '...therapy to treat cancer...'",
        "• Result: Drifts from 'Risk' to 'Cure'.",
        "• Fix: Fusion w/ RM3/SDM anchors this drift."
    ], left=Inches(7), width=Inches(5.8), top=Inches(2.5))

    # 7. RUN 2: SDM
    slide = setup_slide(prs)
    add_title(slide, "Run 2: Structural Validation (SDM)", "Statistical Precision")
    add_card_content(slide, "Beyond Keywords", [
        "• Replaces standard BM25/RM3.",
        "• Logic: Terms appearing as PHRASES are more relevant.",
        "• Zero Parameters: Purely statistical (Markov Random Field).",
        "• Result: Improvements in precision without training."
    ])

    # 8. RUN 3: MONOT5
    slide = setup_slide(prs)
    add_title(slide, "Run 3: Neural Reranking", "Passage-Level Scoring")
    add_card_content(slide, "Deep Scoring", [
        "• Model: MonoT5-3B (Adapated via DuqGen).",
        "• Why DuqGen? learned corpus structure via synthetic queries (Zero-Shot).",
        "• Logic: All Run 1 docs are re-scored.",
        "   - Documents split into passages.",
        "   - MaxP: Doc score = Score of best passage."
    ])

    # 9. HANDLING LONG DOCS
    slide = setup_slide(prs)
    add_title(slide, "Engineering: Handling Long Docs", "MaxP Aggregation")
    add_card_content(slide, "Single GPU Constraint", [
        "• Robust04 docs are long. T5 limit is 512 tokens.",
        "• Method: Sliding Windows.",
        "• Efficiency: We only score the top candidates from Run 1.",
        "• Allows using a 3B parameter model on one RTX 5090."
    ])

    # 10. OPTIMIZATION
    slide = setup_slide(prs)
    add_title(slide, "Optimization & Tuning", "Hyperparameters, not Weights")
    add_card_content(slide, "Refining the Assembly", [
        "• HyDE Prompts: Tested templates ('Write a news article...' wins).",
        "• Fusion Weights: Grid search (0.4 SPLADE / 0.4 Dense / 0.2 Lexical).",
        "• Ablation: Confirmed Dense-Only fails (0.19 MAP) without Fusion (0.30 MAP)."
    ])

    # 11. RESULTS
    slide = setup_slide(prs)
    add_title(slide, "Final Results", "Performance by Run")
    add_card_content(slide, "MAP Scores (Judged)", [
        "Baseline (BM25 + RM3): 0.272",
        "Run 1 (Fusion Baseline): 0.306",
        "Run 2 (SDM Variant): 0.322",
        "Run 3 (Neural Rerank): 0.378  <-- FINAL SUBMISSION",
        "",
        "Total Improvement: +39%"
    ])

    # 12. DISCUSSION
    slide = setup_slide(prs)
    add_title(slide, "Discussion & Takeaway", "Conclusion")
    add_card_content(slide, None, [
        "• Small Data requires assembly, not training.",
        "• Diversity Matters: Fusing Sparse + Dense > Single Model.",
        "• Adaptation Matters: DuqGen (Synthetic) provided the domain knowledge needed for the 3B model to succeed.",
        "• The system is robust, efficient, and statistically safe."
    ])

    # 13. QUESTIONS
    slide = setup_slide(prs)
    tb = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11.3), Inches(2))
    p = tb.text_frame.paragraphs[0]
    p.text = "QUESTIONS?"
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Segoe UI Black'
    p.font.size = Pt(64)
    p.font.color.rgb = RGBColor(255, 255, 255)

    prs.save('Final_Project_Presentation_Extended_v8.pptx')
    print("Done Extended v8")

if __name__ == "__main__":
    create_presentation()
