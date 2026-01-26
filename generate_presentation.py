from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image, ImageDraw
import tempfile
import os
import math

# --- KONFIGURASI DESAIN PREMIUM (Modern Dark Theme) ---
THEME = {
    'bg_dark': (20, 20, 20),           # Almost Black
    'bg_card': (35, 35, 35),           # Dark Gray for Cards
    'text_primary': (255, 255, 255),   # White
    'text_secondary': (180, 180, 180), # Light Gray
    'accent_gradient_start': (0, 120, 215),  # Microsoft Blue
    'accent_gradient_end': (0, 255, 255),    # Cyan
    'accent_warning': (255, 100, 100)        # Soft Red for warnings/failures
}

def create_gradient_img(width, height, start_color, end_color):
    """Create a high-quality diagonal gradient."""
    img = Image.new('RGB', (width, height))
    draw = ImageDraw.Draw(img)
    for x in range(width):
        for y in range(height):
            # Diagonal gradient calculation
            ratio = (x/width + y/height) / 2
            r = int(start_color[0] * (1 - ratio) + end_color[0] * ratio)
            g = int(start_color[1] * (1 - ratio) + end_color[1] * ratio)
            b = int(start_color[2] * (1 - ratio) + end_color[2] * ratio)
            img.putpixel((x, y), (r, g, b))
    return img

def add_gradient_shape(slide, left, top, width, height):
    """Adds a shape filled with the accent gradient (using an image fill workaround)."""
    # 1. Create gradient image
    w_px, h_px = int(width.inches * 96), int(height.inches * 96) # approx DPI
    if w_px < 1 or h_px < 1: return
    
    img = create_gradient_img(w_px, h_px, THEME['accent_gradient_start'], THEME['accent_gradient_end'])
    
    with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as tmp:
        img.save(tmp.name)
        tmp_path = tmp.name
        
    # 2. Add picture as a "shape"
    pic = slide.shapes.add_picture(tmp_path, left, top, width, height)
    
    if os.path.exists(tmp_path): os.unlink(tmp_path)
    return pic

def setup_slide(prs, layout_index=6): # 6 = Blank
    slide = prs.slides.add_slide(prs.slide_layouts[layout_index])
    
    # Dark Background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid() # CRITICAL FIX: Initialize as solid fill
    background.fill.fore_color.rgb = RGBColor(*THEME['bg_dark'])
    background.line.fill.background() # No border
    
    # Decorative Top Bar (Gradient)
    add_gradient_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.15))
    
    return slide

def add_title(slide, text, subtitle=None):
    # Main Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
    p = title_box.text_frame.paragraphs[0]
    p.text = text.upper()
    p.font.name = 'Segoe UI Black' # Heavier font
    p.font.size = Pt(36)
    p.font.color.rgb = RGBColor(*THEME['text_primary'])
    # p.space_after = Pt(6) # Removing incompatible properties if any
    
    # Accent Line under title
    add_gradient_shape(slide, Inches(0.5), Inches(1.3), Inches(2), Inches(0.05))

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(12), Inches(0.8))
        p = sub_box.text_frame.paragraphs[0]
        p.text = subtitle
        p.font.name = 'Segoe UI Light'
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(*THEME['text_secondary'])

def add_card_content(slide, title, bullets, left=Inches(0.5), top=Inches(2.2), width=Inches(12.3)):
    """Adds a 'card' style content box."""
    # Card Background shape
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, Inches(4.8))
    card.fill.solid() # CRITICAL FIX
    card.fill.fore_color.rgb = RGBColor(*THEME['bg_card'])
    card.line.fill.background()
    card.shadow.inherit = False # Clean look
    
    # Content Title (inside card)
    if title:
        tb = slide.shapes.add_textbox(left + Inches(0.3), top + Inches(0.2), width - Inches(0.6), Inches(0.8))
        p = tb.text_frame.paragraphs[0]
        p.text = title
        p.font.name = 'Segoe UI Semibold'
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(*THEME['accent_gradient_end']) # Cyan text
        
    # Bullets
    border_offset = Inches(0.8) if title else Inches(0.3)
    tb = slide.shapes.add_textbox(left + Inches(0.3), top + border_offset, width - Inches(0.6), Inches(3.8))
    tf = tb.text_frame
    tf.word_wrap = True
    
    for line in bullets:
        p = tf.add_paragraph()
        p.text = line
        p.font.name = 'Segoe UI'
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(*THEME['text_primary'])
        p.space_before = Pt(12)
        p.space_after = Pt(12)

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # --- SLIDE 1: HERO TITLE ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Full Gradient Background
    w_px, h_px = 1920, 1080
    bg_img = Image.new('RGB', (w_px, h_px))
    draw = ImageDraw.Draw(bg_img)
    # Dark Radial-like Gradient manually
    # Simple linear diagonal for now, but dark to blue
    for x in range(w_px):
        for y in range(h_px):
            ratio = (x + y) / (w_px + h_px)
            # Fade from Black (0,0,0) to deep blue (0, 20, 60)
            r = int(0 + 0 * ratio)
            g = int(0 + 20 * ratio)
            b = int(10 + 60 * ratio)
            bg_img.putpixel((x, y), (r, g, b))
            
    with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as tmp:
        bg_img.save(tmp.name)
        tmp_path = tmp.name
    slide.shapes.add_picture(tmp_path, 0, 0, prs.slide_width, prs.slide_height)
    if os.path.exists(tmp_path): os.unlink(tmp_path)
    
    # Center Title
    tb = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.3), Inches(2.5))
    p = tb.text_frame.paragraphs[0]
    p.text = "MAXIMIZING RETRIEVAL\nEFFECTIVENESS ON ROBUST04"
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Segoe UI Black'
    p.font.size = Pt(54)
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Subtitle
    tb = slide.shapes.add_textbox(Inches(2), Inches(5), Inches(9.3), Inches(1))
    p = tb.text_frame.paragraphs[0]
    p.text = "Final Project - Part A | Text Retrieval & Search Engines | 2026"
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Segoe UI Light'
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(0, 255, 255) # Cyan accent
    
    # --- SLIDE 2: CHALLENGE ---
    slide = setup_slide(prs)
    add_title(slide, "The Challenge", "Why is Robust04 difficult?")
    
    add_card_content(slide, "Dataset Characteristics", [
        "• 250 Queries with highly variable difficulty.",
        "• Documents are extremely long (truncation risks).",
        "• Hard relevance judgments requiring high precision.",
        "• Baseline Goal: Beat BM25+RM3 (MAP 0.2719)."
    ])
    
    # --- SLIDE 3: PIPELINE ---
    slide = setup_slide(prs)
    add_title(slide, "Our Architecture", "A 3-Stage Hybrid System")
    
    # Draw simple flow diagram using shapes
    # Box 1
    add_gradient_shape(slide, Inches(1), Inches(3), Inches(3), Inches(2))
    tb = slide.shapes.add_textbox(Inches(1.2), Inches(3.5), Inches(2.6), Inches(1))
    p = tb.text_frame.paragraphs[0]
    p.text = "Stage 1\nHybrid Retrieval"
    p.alignment = PP_ALIGN.CENTER
    p.font.color.rgb = RGBColor(255,255,255)
    p.font.bold = True
    
    # Arrow
    sh = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(4.2), Inches(3.8), Inches(1), Inches(0.4))
    sh.fill.solid(); sh.fill.fore_color.rgb = RGBColor(100,100,100); sh.line.fill.background()

    # Box 2
    add_gradient_shape(slide, Inches(5.4), Inches(3), Inches(3), Inches(2))
    tb = slide.shapes.add_textbox(Inches(5.6), Inches(3.5), Inches(2.6), Inches(1))
    p = tb.text_frame.paragraphs[0]
    p.text = "Stage 2\nWeighted Fusion"
    p.alignment = PP_ALIGN.CENTER
    p.font.color.rgb = RGBColor(255,255,255)
    p.font.bold = True
    
    # Arrow
    sh = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(8.6), Inches(3.8), Inches(1), Inches(0.4))
    sh.fill.solid(); sh.fill.fore_color.rgb = RGBColor(100,100,100); sh.line.fill.background()

    # Box 3 (Highlight)
    add_gradient_shape(slide, Inches(9.8), Inches(2.8), Inches(3.2), Inches(2.4)) # Larger
    tb = slide.shapes.add_textbox(Inches(10), Inches(3.5), Inches(2.8), Inches(1))
    p = tb.text_frame.paragraphs[0]
    p.text = "Stage 3\nNeural Reranking\n(MonoT5)"
    p.alignment = PP_ALIGN.CENTER
    p.font.color.rgb = RGBColor(255,255,255)
    p.font.bold = True
    p.font.size = Pt(24)

    # --- SLIDE 4: METHOD 1 ---
    slide = setup_slide(prs)
    add_title(slide, "Dense-HyDE Hybrid", "Solving the 'Hallucination vs Drift' Dilemma")
    
    # Split layout
    add_card_content(slide, "The Problem", [
        "• HyDE generates fake document passages.",
        "• Dense Retrievers (BGE) love this semantic context.",
        "• Lexical Retrievers (BM25) get confused by fake terms."
    ], width=Inches(6), top=Inches(2.5))
    
    add_card_content(slide, "The Solution", [
        "• Split the Query processing:",
        "• Lexical Models → See Original Query", 
        "• Dense Models → See HyDE Passage",
        "• Result: +12% improvement over baseline."
    ], left=Inches(6.8), width=Inches(6), top=Inches(2.5))

    # --- SLIDE 5: METHOD 2 (SDM Validation) ---
    slide = setup_slide(prs)
    add_title(slide, "Validation: Sequential Dependence (SDM)", "Why Structure Matters (Experiment)")
    
    add_card_content(slide, "Modeling Term Dependencies", [
        "• Standard Bag-of-Words (BM25) ignores phrase structure.",
        "• SDM (Markov Random Field) captures local dependencies.",
        "• Weights: Term (0.75), Ordered (0.10), Unordered (0.15).",
        "• Outcome: Validated stability (+0.001 MAP), but neural fusion (Run 3) proved superior without it."
    ])

    # --- SLIDE 6: METHOD 3 (MonoT5) ---
    slide = setup_slide(prs)
    add_title(slide, "Neural Reranking: MonoT5", "The Performance Driver")
    
    add_card_content(slide, None, [
        "Model: monot5-3b-robust04 (3 Billion Parameters)",
        "Technique: MaxP Aggregation (Max Passage)",
        "• Motivation: Robust04 docs are long; relevance is sparse.",
        "• Method: Score all sliding windows, take MAX score.",
        "• Theoretical Basis: Aligns with 'Birch' (Yilmaz et al., 2019) and 'PARADE-Max' (Li et al., 2020) findings.",
        "• Result: +39% improvement over baseline."
    ])


    # --- SLIDE 7: RESULTS ---
    slide = setup_slide(prs)
    add_title(slide, "Results Summary", "Total Relative Improvement: +39%")
    
    # Bar Chart Visual (Shapes)
    # Baseline
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2), Inches(5), Inches(2), Inches(2))
    s.fill.solid(); s.fill.fore_color.rgb = RGBColor(100,100,100)
    
    tb = slide.shapes.add_textbox(Inches(2), Inches(4.5), Inches(2), Inches(0.5))
    tb.text_frame.text = "0.272"
    tb.text_frame.paragraphs[0].font.color.rgb = RGBColor(200,200,200); tb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    tb = slide.shapes.add_textbox(Inches(2), Inches(7.1), Inches(2), Inches(0.5))
    tb.text_frame.text = "Baseline\n(BM25)"; tb.text_frame.paragraphs[0].font.color.rgb = RGBColor(200,200,200); tb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Fusion
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5), Inches(4.5), Inches(2), Inches(2.5))
    s.fill.solid(); s.fill.fore_color.rgb = RGBColor(*THEME['accent_gradient_start'])
    
    tb = slide.shapes.add_textbox(Inches(5), Inches(4), Inches(2), Inches(0.5))
    tb.text_frame.text = "0.306"
    tb.text_frame.paragraphs[0].font.color.rgb = RGBColor(255,255,255); tb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    tb = slide.shapes.add_textbox(Inches(5), Inches(7.1), Inches(2), Inches(0.5))
    tb.text_frame.text = "Fusion\n(Dense-HyDE)"; tb.text_frame.paragraphs[0].font.color.rgb = RGBColor(200,200,200); tb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # MonoT5
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8), Inches(3), Inches(2), Inches(4))
    s.fill.solid(); s.fill.fore_color.rgb = RGBColor(*THEME['accent_gradient_end'])
    
    tb = slide.shapes.add_textbox(Inches(8), Inches(2.5), Inches(2), Inches(0.5))
    tb.text_frame.text = "0.378"
    tb.text_frame.paragraphs[0].font.color.rgb = RGBColor(0,255,255); tb.text_frame.paragraphs[0].font.bold=True; tb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    tb = slide.shapes.add_textbox(Inches(8), Inches(7.1), Inches(2), Inches(0.5))
    tb.text_frame.text = "Final\n(MonoT5 MaxP)"; tb.text_frame.paragraphs[0].font.color.rgb = RGBColor(200,200,200); tb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # --- SLIDE 8: REFERENCES ---
    slide = setup_slide(prs)
    add_title(slide, "Key References", "Academic Foundation")
    
    add_card_content(slide, None, [
        "1. Nogueira et al. (2020). Document Ranking with a Pretrained Sequence-to-Sequence Model. (MonoT5)",
        "2. Yilmaz et al. (2019). Cross-Domain Modeling of Sentence-Level Evidence for Document Retrieval. (Birch / MaxP)",
        "3. Li et al. (2020). PARADE: Passage Representation Aggregation for Document Reranking.",
        "4. Gao et al. (2022). Precise Zero-Shot Dense Retrieval without Relevance Labels. (HyDE)"
    ])

    # --- SLIDE 9: CONCLUSION ---
    slide = setup_slide(prs)
    
    # Centered Big Finish
    tb = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11.3), Inches(2))
    p = tb.text_frame.paragraphs[0]
    p.text = "FINAL MAP: 0.378"
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Segoe UI Black'
    p.font.size = Pt(64)
    p.font.color.rgb = RGBColor(0, 255, 255)
    
    tb = slide.shapes.add_textbox(Inches(3), Inches(5), Inches(7.3), Inches(1))
    p = tb.text_frame.paragraphs[0]
    p.text = "Thank You"
    p.alignment = PP_ALIGN.CENTER
    p.font.color.rgb = RGBColor(200, 200, 200)

    prs.save('Final_Project_Presentation_Premium_v3.pptx')
    print("Done")

if __name__ == "__main__":
    create_presentation()
